"""
מצגת מעוף – Tech-Lead Israel
Professional Pitch Deck
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# ─── Color Palette ───
NAVY      = RGBColor(0x1E, 0x27, 0x61)
DARK_NAVY = RGBColor(0x14, 0x1B, 0x41)
TEAL      = RGBColor(0x00, 0xA8, 0x96)
GOLD      = RGBColor(0xF9, 0xC7, 0x4F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
CORAL     = RGBColor(0xF9, 0x61, 0x67)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
LIGHT_TEAL= RGBColor(0xCC, 0xFB, 0xF1)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, x, y, w, h, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.RIGHT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_multiline_box(slide, x, y, w, h, lines, font_size=14, color=WHITE, alignment=PP_ALIGN.RIGHT, line_spacing=1.5, font_name="Arial"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)

        if isinstance(line_data, dict):
            run = p.add_run()
            run.text = line_data.get("text", "")
            run.font.size = Pt(line_data.get("size", font_size))
            run.font.color.rgb = line_data.get("color", color)
            run.font.bold = line_data.get("bold", False)
            run.font.name = font_name
        else:
            run = p.add_run()
            run.text = str(line_data)
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
    return txBox

def add_card(slide, x, y, w, h, title, body_lines, accent_color=TEAL):
    """Add a card with accent bar on right side"""
    # Card background
    card = add_rect(slide, x, y, w, h, WHITE)
    # Accent bar (right side for RTL)
    add_rect(slide, Emu(x.emu + w.emu - Inches(0.06).emu), y, Inches(0.06), h, accent_color)
    # Title
    add_text_box(slide, Emu(x.emu + Inches(0.2).emu), y, Emu(w.emu - Inches(0.5).emu), Inches(0.5),
                 title, font_size=16, color=DARK_TEXT, bold=True)
    # Body
    body_y = Emu(y.emu + Inches(0.45).emu)
    for i, line in enumerate(body_lines):
        add_text_box(slide, Emu(x.emu + Inches(0.2).emu), Emu(body_y.emu + Inches(i * 0.3).emu),
                     Emu(w.emu - Inches(0.5).emu), Inches(0.3),
                     line, font_size=12, color=GRAY_TEXT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_NAVY)

# Decorative shapes
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)
add_rect(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GOLD)

# Side accent
add_rect(slide, Inches(12.5), Inches(0.5), Inches(0.06), Inches(6.5), TEAL)

# Main title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "מעוף", font_size=72, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Tech-Lead Israel", font_size=36, color=TEAL, bold=True)

add_multiline_box(slide, Inches(1), Inches(4.0), Inches(10), Inches(1.5), [
    {"text": "מערכת AI לתמיכת החלטה – גיוס מומחי STEM למערכת החינוך", "size": 20, "color": WHITE},
    {"text": "במודל קריירה מעגלית 2+2", "size": 20, "color": GOLD},
])

add_text_box(slide, Inches(1), Inches(6.0), Inches(10), Inches(0.5),
             "Confidential | 2026", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "הבעיה – שני צדדים של אותו מטבע", font_size=36, color=NAVY, bold=True)

# Left crisis - Education
add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), WHITE)
add_rect(slide, Inches(12.54), Inches(1.5), Inches(0.06), Inches(5.2), CORAL)

add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.6),
             "משבר מורי STEM", font_size=24, color=CORAL, bold=True)

crisis_edu = [
    {"text": "3,500", "size": 48, "color": CORAL, "bold": True},
    {"text": "משרות הוראה פנויות במקצועות STEM", "size": 16, "color": DARK_TEXT},
    {"text": "", "size": 8, "color": DARK_TEXT},
    {"text": "22% ירידה בגיוס מורים חדשים", "size": 16, "color": GRAY_TEXT},
    {"text": "60% מהמורים עוזבים תוך 5 שנים", "size": 16, "color": GRAY_TEXT},
    {"text": "פריפריה מפסידה הכי הרבה", "size": 16, "color": GRAY_TEXT},
]
add_multiline_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(3.5), crisis_edu)

# Right crisis - Hi-tech
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), WHITE)
add_rect(slide, Inches(6.24), Inches(1.5), Inches(0.06), Inches(5.2), TEAL)

add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4.8), Inches(0.6),
             "משבר גיוס ג'וניורים", font_size=24, color=TEAL, bold=True)

crisis_tech = [
    {"text": "38-43K", "size": 48, "color": TEAL, "bold": True},
    {"text": "עלות גיוס ממוצעת לעובד ג'וניור (בש\"ח)", "size": 16, "color": DARK_TEXT},
    {"text": "", "size": 8, "color": DARK_TEXT},
    {"text": "3 חודשי אחריות בלבד מחברות השמה", "size": 16, "color": GRAY_TEXT},
    {"text": "40% נטישה בשנה הראשונה", "size": 16, "color": GRAY_TEXT},
    {"text": "חוסר Soft Skills אצל בוגרי בוטקמפ", "size": 16, "color": GRAY_TEXT},
]
add_multiline_box(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.5), crisis_tech)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: Target Audience (UPDATED - not just combat soldiers)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "קהל היעד – חיילים משוחררים עם פוטנציאל STEM", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
             "לא רק לוחמים – כל חייל עם רקע טכנולוגי או מקצועי", font_size=18, color=GOLD)

# Row of audience cards
audiences = [
    ("לוחמים משוחררים", "רקע מבצעי, מנהיגות\nעמידה בלחץ, עבודת צוות"),
    ("תקשוב ומערכות מידע", "ניסיון IT, רשתות, סייבר\nידע טכני מעשי"),
    ("חיל האוויר – טכני", "אלקטרוניקה, תעופה\nתחזוקה, מערכות מורכבות"),
    ("מודיעין ולוט\"ם", "ניתוח נתונים, AI/ML\nתכנות, מחקר"),
    ("ממר\"ם ומנה\"ט", "פיתוח תוכנה, DevOps\nפרויקטים טכנולוגיים"),
]

card_w = Inches(2.2)
card_h = Inches(3.8)
start_x = Inches(0.5)
gap = Inches(0.28)

for i, (title, desc) in enumerate(audiences):
    cx = Emu(start_x.emu + i * (card_w.emu + gap.emu))
    cy = Inches(2.0)

    # Card bg
    card = add_rect(slide, cx, cy, card_w, card_h, RGBColor(0x1B, 0x25, 0x55))
    # Top accent
    add_rect(slide, cx, cy, card_w, Inches(0.06), TEAL if i > 0 else GOLD)

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(cx.emu + card_w.emu // 2 - Inches(0.4).emu), Emu(cy.emu + Inches(0.4).emu),
        Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL if i > 0 else GOLD
    circle.line.fill.background()

    # Icon text (emoji placeholder)
    icons = ["⚔", "🖥", "✈", "📊", "💻"]
    add_text_box(slide, Emu(cx.emu + card_w.emu // 2 - Inches(0.35).emu),
                 Emu(cy.emu + Inches(0.5).emu), Inches(0.7), Inches(0.6),
                 icons[i], font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Emu(cx.emu + Inches(0.15).emu), Emu(cy.emu + Inches(1.5).emu),
                 Emu(card_w.emu - Inches(0.3).emu), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, Emu(cx.emu + Inches(0.15).emu),
                     Emu(cy.emu + Inches(2.1).emu + Inches(j * 0.35).emu),
                     Emu(card_w.emu - Inches(0.3).emu), Inches(0.35),
                     line, font_size=12, color=RGBColor(0xA0, 0xAE, 0xC0), alignment=PP_ALIGN.CENTER)

# Bottom stat
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
             "פוטנציאל: מעל 15,000 חיילים משוחררים בשנה עם רקע טכנולוגי רלוונטי",
             font_size=16, color=GOLD, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: The Solution - 2+2 Model
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "הפתרון – מודל 2+2", font_size=36, color=NAVY, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.5),
             "הכשרה ← 2 שנות הוראת STEM ← 2 שנות הייטק מובטחות", font_size=18, color=TEAL, bold=True)

# Phase cards - 4 steps
phases = [
    ("שלב 0 – הכשרה", "קורס AI & Data Science", "בנאיה קולג'\nמימון קרן לחיילים\nמשוחררים", TEAL),
    ("שלב 1 – הוראת STEM", "2 שנות תרומה", "סייבר, פייתון, פיזיקה\nחוזה אישי\nשכר ~15,000 ש\"ח", GOLD),
    ("שלב 2 – תעשיית הייטק", "2 שנות קריירה מובטחות", "Microsoft, Elbit\nCheckpoint\nחברות מאמצות", GREEN),
    ("ערך מוסף", "Seniority Leap", "2 שנות הוראה =\n18 חודשי ניסיון מקצועי\nSoft Skills + ניהול", CORAL),
]

phase_w = Inches(2.8)
phase_h = Inches(4.0)
phase_gap = Inches(0.3)
phase_start = Inches(0.6)

for i, (title, subtitle, desc, accent) in enumerate(phases):
    px = Emu(phase_start.emu + i * (phase_w.emu + phase_gap.emu))
    py = Inches(2.0)

    # Card
    add_rect(slide, px, py, phase_w, phase_h, WHITE)
    # Top accent bar
    add_rect(slide, px, py, phase_w, Inches(0.08), accent)

    # Phase number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(px.emu + phase_w.emu // 2 - Inches(0.35).emu), Emu(py.emu + Inches(0.3).emu),
        Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()

    num_text = str(i) if i < 3 else "+"
    add_text_box(slide, Emu(px.emu + phase_w.emu // 2 - Inches(0.3).emu),
                 Emu(py.emu + Inches(0.4).emu), Inches(0.6), Inches(0.5),
                 num_text, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Emu(px.emu + Inches(0.2).emu), Emu(py.emu + Inches(1.2).emu),
                 Emu(phase_w.emu - Inches(0.4).emu), Inches(0.5),
                 title, font_size=16, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Emu(px.emu + Inches(0.2).emu), Emu(py.emu + Inches(1.7).emu),
                 Emu(phase_w.emu - Inches(0.4).emu), Inches(0.4),
                 subtitle, font_size=13, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    # Description lines
    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, Emu(px.emu + Inches(0.2).emu),
                     Emu(py.emu + Inches(2.3).emu + Inches(j * 0.35).emu),
                     Emu(phase_w.emu - Inches(0.4).emu), Inches(0.35),
                     line, font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Arrow indicators between cards
for i in range(3):
    ax = Emu(phase_start.emu + (i + 1) * (phase_w.emu + phase_gap.emu) - phase_gap.emu // 2 - Inches(0.15).emu)
    add_text_box(slide, ax, Inches(3.8), Inches(0.3), Inches(0.4),
                 "←", font_size=24, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Dual Scoring Engine
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
             "Dual Scoring Engine – מנוע דירוג כפול", font_size=32, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4),
             "כל מועמד מקבל שני ציונים נפרדים – חובה לפחות 60 בכל אחד", font_size=16, color=GOLD)

# Score A - School Match (right side)
add_rect(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.0), RGBColor(0x1B, 0x25, 0x55))
add_rect(slide, Inches(12.74), Inches(1.7), Inches(0.06), Inches(5.0), TEAL)

add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.2), Inches(0.5),
             "ציון A: התאמה לבית ספר", font_size=22, color=TEAL, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.4), Inches(5.2), Inches(0.4),
             "School Match – 100 נקודות", font_size=14, color=GOLD)

score_a = [
    ("כושר הוראה", "35%", "מבחן ELI5 + תקשורת + הדרכה"),
    ("התאמה מקצועית", "25%", "מומחיות vs צורך ביה\"ס"),
    ("שימור", "25%", "מרחק + משפחה + מחויבות"),
    ("Impact", "15%", "פריפריה + טיפוח + דחיפות"),
]

for i, (name, weight, desc) in enumerate(score_a):
    row_y = Emu(Inches(3.0).emu + Inches(i * 0.65).emu)
    # Weight badge
    badge = add_rect(slide, Inches(11.2), row_y, Inches(0.9), Inches(0.4), TEAL)
    add_text_box(slide, Inches(11.2), row_y, Inches(0.9), Inches(0.4),
                 weight, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide, Inches(9.5), row_y, Inches(1.6), Inches(0.4),
                 name, font_size=14, color=WHITE, bold=True)
    # Desc
    add_text_box(slide, Inches(7.2), row_y, Inches(2.2), Inches(0.4),
                 desc, font_size=11, color=RGBColor(0xA0, 0xAE, 0xC0))

# Score B - Company Match (left side)
add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.0), RGBColor(0x1B, 0x25, 0x55))
add_rect(slide, Inches(6.24), Inches(1.7), Inches(0.06), Inches(5.0), GOLD)

add_text_box(slide, Inches(0.9), Inches(1.9), Inches(5.0), Inches(0.5),
             "ציון B: התאמה לחברה", font_size=22, color=GOLD, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.4), Inches(5.0), Inches(0.4),
             "Company Match – 100 נקודות", font_size=14, color=TEAL)

score_b = [
    ("מיומנות טכנית", "35%", "Tech-Stack + מבחן + רקע צבאי"),
    ("פוטנציאל צמיחה", "25%", "קצב למידה + פסיכומטרי"),
    ("כישורים רכים", "25%", "ELI5 + מנהיגות + צוות"),
    ("התאמה תרבותית", "15%", "גודל חברה + סגנון עבודה"),
]

for i, (name, weight, desc) in enumerate(score_b):
    row_y = Emu(Inches(3.0).emu + Inches(i * 0.65).emu)
    badge = add_rect(slide, Inches(4.8), row_y, Inches(0.9), Inches(0.4), GOLD)
    add_text_box(slide, Inches(4.8), row_y, Inches(0.9), Inches(0.4),
                 weight, font_size=14, color=DARK_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.0), row_y, Inches(1.7), Inches(0.4),
                 name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.9), row_y, Inches(2.0), Inches(0.4),
                 desc, font_size=11, color=RGBColor(0xA0, 0xAE, 0xC0))

# Bottom: Placement Score
add_rect(slide, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.7), TEAL)
add_text_box(slide, Inches(3.5), Inches(5.65), Inches(6.3), Inches(0.6),
             "ציון שיבוץ = שילוב משוקלל A + B  |  סף מינימום: 60 בכל ציון",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Three-Sided Marketplace
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "שוק תלת-צדדי – ערך לכל השחקנים", font_size=36, color=NAVY, bold=True)

stakeholders = [
    ("חיילים משוחררים", TEAL, [
        "הכשרה מקצועית במימון מלא",
        "שכר 15K ש\"ח בזמן ההוראה",
        "2 שנות הייטק מובטחות",
        "Seniority Leap: +18 חודשי ניסיון",
        "מסלול קריירה ברור ויציב",
    ]),
    ("בתי ספר / משרד החינוך", GOLD, [
        "מורי STEM איכותיים ומתחלפים",
        "מימון 100% מתקציב גפ\"ן",
        "מילוי משרות בפריפריה",
        "עדכון תוכניות לימוד טכנולוגיות",
        "AI Matching – שיבוץ אופטימלי",
    ]),
    ("חברות הייטק מאמצות", CORAL, [
        "עובדים מסוננים עם 2 שנות ניסיון",
        "חיסכון 200K ש\"ח vs גיוס רגיל",
        "אחריות 24 חודש (לא 3)",
        "Soft Skills מוכחים מהשטח",
        "Pipeline קבוע של טאלנטים",
    ]),
]

card_w = Inches(3.8)
card_h = Inches(4.8)
gap = Inches(0.35)
start_x = Inches(0.65)

for i, (title, accent, items) in enumerate(stakeholders):
    cx = Emu(start_x.emu + i * (card_w.emu + gap.emu))
    cy = Inches(1.8)

    add_rect(slide, cx, cy, card_w, card_h, WHITE)
    add_rect(slide, cx, cy, card_w, Inches(0.08), accent)

    # Circle icon
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(cx.emu + card_w.emu // 2 - Inches(0.35).emu), Emu(cy.emu + Inches(0.3).emu),
        Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()

    icons = ["🎖", "🏫", "🏢"]
    add_text_box(slide, Emu(cx.emu + card_w.emu // 2 - Inches(0.3).emu),
                 Emu(cy.emu + Inches(0.38).emu), Inches(0.6), Inches(0.5),
                 icons[i], font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Emu(cx.emu + Inches(0.3).emu), Emu(cy.emu + Inches(1.2).emu),
                 Emu(card_w.emu - Inches(0.6).emu), Inches(0.5),
                 title, font_size=18, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text_box(slide, Emu(cx.emu + Inches(0.3).emu),
                     Emu(cy.emu + Inches(1.9).emu + Inches(j * 0.45).emu),
                     Emu(card_w.emu - Inches(0.6).emu), Inches(0.4),
                     "◂  " + item, font_size=13, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: Unit Economics
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
             "Unit Economics – הכלכלה עובדת", font_size=36, color=WHITE, bold=True)

# Big stats row
stats = [
    ("83K ₪", "הכנסה למשתתף", TEAL),
    ("11K ₪", "עלות למשתתף", CORAL),
    ("87%", "מרווח גולמי", GREEN),
    ("30", "Break-Even (חיילים)", GOLD),
]

stat_w = Inches(2.8)
for i, (num, label, clr) in enumerate(stats):
    sx = Emu(Inches(0.6).emu + i * (stat_w.emu + Inches(0.25).emu))
    sy = Inches(1.5)

    add_rect(slide, sx, sy, stat_w, Inches(1.8), RGBColor(0x1B, 0x25, 0x55))
    add_rect(slide, sx, sy, stat_w, Inches(0.06), clr)

    add_text_box(slide, sx, Emu(sy.emu + Inches(0.3).emu), stat_w, Inches(0.8),
                 num, font_size=40, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, sx, Emu(sy.emu + Inches(1.1).emu), stat_w, Inches(0.4),
                 label, font_size=14, color=RGBColor(0xA0, 0xAE, 0xC0), alignment=PP_ALIGN.CENTER)

# Revenue breakdown
add_text_box(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(0.5),
             "מקורות הכנסה", font_size=20, color=TEAL, bold=True)

revenues = [
    ("עמלת השמה לחברה", "35,000 ₪", "חד-פעמי בתחילת שלב הייטק"),
    ("דמי ניהול SaaS", "18,000 ₪/שנה", "ממשרד החינוך על הפלטפורמה"),
    ("מענק חתימה (פריפריה)", "עד 45,000 ₪", "תקציב ממשלתי קיים – עובר דרך מעוף"),
    ("סבסוד דיור (ביקוש)", "2,200 ₪/חודש", "אזורי ת\"א – מגפ\"ן"),
]

for i, (name, amount, desc) in enumerate(revenues):
    ry = Emu(Inches(4.3).emu + Inches(i * 0.65).emu)
    add_rect(slide, Inches(0.8), ry, Inches(5.5), Inches(0.55), RGBColor(0x1B, 0x25, 0x55))
    add_text_box(slide, Inches(4.3), ry, Inches(2.0), Inches(0.55),
                 amount, font_size=14, color=GOLD, bold=True)
    add_text_box(slide, Inches(2.5), ry, Inches(1.7), Inches(0.55),
                 name, font_size=13, color=WHITE)
    add_text_box(slide, Inches(0.9), ry, Inches(1.5), Inches(0.55),
                 desc, font_size=10, color=RGBColor(0xA0, 0xAE, 0xC0))

# Comparison table
add_text_box(slide, Inches(7.0), Inches(3.7), Inches(5.5), Inches(0.5),
             "מעוף vs גיוס רגיל", font_size=20, color=GOLD, bold=True)

add_rect(slide, Inches(7.0), Inches(4.3), Inches(5.8), Inches(0.5), TEAL)
add_text_box(slide, Inches(10.5), Inches(4.3), Inches(2.3), Inches(0.5),
             "מדד", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.8), Inches(4.3), Inches(1.7), Inches(0.5),
             "מעוף", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.0), Inches(4.3), Inches(1.8), Inches(0.5),
             "גיוס רגיל", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

comparisons = [
    ("עלות גיוס", "35K ₪", "38-43K ₪"),
    ("אחריות", "24 חודש", "3 חודש"),
    ("Soft Skills", "מוכחים (2 שנות הוראה)", "לא ידוע"),
    ("נטישה שנה 1", "צפי: ~10%", "~40%"),
]

for i, (metric, maof, regular) in enumerate(comparisons):
    ry = Emu(Inches(4.85).emu + Inches(i * 0.5).emu)
    bg_clr = RGBColor(0x1B, 0x25, 0x55) if i % 2 == 0 else RGBColor(0x16, 0x1F, 0x48)
    add_rect(slide, Inches(7.0), ry, Inches(5.8), Inches(0.45), bg_clr)
    add_text_box(slide, Inches(10.5), ry, Inches(2.3), Inches(0.45),
                 metric, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.8), ry, Inches(1.7), Inches(0.45),
                 maof, font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), ry, Inches(1.8), Inches(0.45),
                 regular, font_size=12, color=CORAL, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: AI Technology Platform
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "הפלטפורמה הטכנולוגית", font_size=36, color=NAVY, bold=True)

# Platform components
components = [
    ("מבחן ELI5", "הסבר מושגים מורכבים\nבשפה פשוטה – 3 סבבים:\nכתיבה, וידאו, צ'אט AI", TEAL),
    ("Dual Scoring Engine", "אלגוריתם AI לדירוג\nוהתאמה אוטומטית\nלבי\"ס + חברה בו-זמנית", GOLD),
    ("פורטל חיילים", "דשבורד אישי, מסלול 2+2\nהתאמות, סימולטור שכר\nמעקב משימות", CORAL),
    ("פורטל משרד החינוך", "מפת צרכים ארצית\n247 מועמדים, AI Insights\nאישורי שיבוץ", GREEN),
    ("פורטל חברות", "ניהול עובדים, מעקב\nROI, תכנון מעבר\nמ-הוראה להייטק", NAVY),
    ("Feedback Loop", "השוואת חיזוי vs ביצועים\nעדכון משקלות אוטומטי\nשיפור מתמיד", RGBColor(0x7C, 0x3A, 0xED)),
]

comp_w = Inches(3.8)
comp_h = Inches(2.0)

for i, (title, desc, accent) in enumerate(components):
    row = i // 3
    col = i % 3
    cx = Emu(Inches(0.65).emu + col * (comp_w.emu + Inches(0.35).emu))
    cy = Emu(Inches(1.6).emu + row * (comp_h.emu + Inches(0.4).emu))

    add_rect(slide, cx, cy, comp_w, comp_h, WHITE)
    # Right accent bar
    add_rect(slide, Emu(cx.emu + comp_w.emu - Inches(0.06).emu), cy, Inches(0.06), comp_h, accent)

    add_text_box(slide, Emu(cx.emu + Inches(0.25).emu), Emu(cy.emu + Inches(0.15).emu),
                 Emu(comp_w.emu - Inches(0.6).emu), Inches(0.4),
                 title, font_size=16, color=accent, bold=True)

    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, Emu(cx.emu + Inches(0.25).emu),
                     Emu(cy.emu + Inches(0.6).emu + Inches(j * 0.35).emu),
                     Emu(comp_w.emu - Inches(0.6).emu), Inches(0.35),
                     line, font_size=12, color=GRAY_TEXT)

# Bottom bar
add_rect(slide, Inches(0.65), Inches(6.2), Inches(12.0), Inches(0.06), TEAL)
add_text_box(slide, Inches(0.65), Inches(6.4), Inches(12.0), Inches(0.5),
             "כל הפורטלים פועלים – גרסת Demo מלאה עם דשבורדים, גרפים, וסימולציות",
             font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: Budget & Ask
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
             "תקציב ובקשת מימון", font_size=36, color=WHITE, bold=True)

# Budget breakdown
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
             "פירוט תקציב – 12 חודשים", font_size=20, color=TEAL, bold=True)

budget_items = [
    ("פיתוח MVP", "90,000 ₪", "38%"),
    ("תשתית ענן ו-AI", "35,000 ₪", "15%"),
    ("עלויות משפט + רגולציה", "25,000 ₪", "11%"),
    ("שיווק ואקוויזיציה", "30,000 ₪", "13%"),
    ("פיילוט שטח", "25,000 ₪", "11%"),
    ("הוצ' כלליות + רזרבה", "30,000 ₪", "13%"),
]

for i, (item, amount, pct) in enumerate(budget_items):
    ry = Emu(Inches(1.9).emu + Inches(i * 0.6).emu)
    bg_clr = RGBColor(0x1B, 0x25, 0x55) if i % 2 == 0 else RGBColor(0x16, 0x1F, 0x48)
    add_rect(slide, Inches(7.0), ry, Inches(5.8), Inches(0.5), bg_clr)

    # Progress bar
    bar_w = float(pct.replace("%", "")) / 100 * Inches(1.2).emu
    add_rect(slide, Inches(7.2), Emu(ry.emu + Inches(0.15).emu), Emu(int(bar_w)), Inches(0.2), TEAL)

    add_text_box(slide, Inches(8.5), ry, Inches(2.0), Inches(0.5),
                 item, font_size=13, color=WHITE)
    add_text_box(slide, Inches(10.8), ry, Inches(1.2), Inches(0.5),
                 amount, font_size=13, color=GOLD, bold=True)
    add_text_box(slide, Inches(12.1), ry, Inches(0.7), Inches(0.5),
                 pct, font_size=11, color=GRAY_TEXT)

# Total
add_rect(slide, Inches(7.0), Emu(Inches(1.9).emu + Inches(6 * 0.6).emu), Inches(5.8), Inches(0.5), TEAL)
add_text_box(slide, Inches(8.5), Emu(Inches(1.9).emu + Inches(6 * 0.6).emu), Inches(2.0), Inches(0.5),
             "סה\"כ", font_size=14, color=WHITE, bold=True)
add_text_box(slide, Inches(10.8), Emu(Inches(1.9).emu + Inches(6 * 0.6).emu), Inches(1.2), Inches(0.5),
             "235,000 ₪", font_size=14, color=WHITE, bold=True)

# Ask panel (left)
add_rect(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.3), RGBColor(0x1B, 0x25, 0x55))
add_rect(slide, Inches(6.44), Inches(1.3), Inches(0.06), Inches(5.3), GOLD)

add_text_box(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.5),
             "הבקשה", font_size=24, color=GOLD, bold=True)

# Big number
add_text_box(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(1.2),
             "200,000 ₪", font_size=56, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.9), Inches(3.5), Inches(5.2), Inches(0.5),
             "מענק רשות החדשנות – מסלול תנופה (85%)", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

ask_details = [
    {"text": "", "size": 10, "color": WHITE},
    {"text": "השקעה עצמית: 35,000 ₪ (15%)", "size": 16, "color": TEAL},
    {"text": "", "size": 10, "color": WHITE},
    {"text": "יעד: MVP עובד + פיילוט 30 חיילים", "size": 16, "color": WHITE},
    {"text": "לו\"ז: 12 חודשים", "size": 16, "color": WHITE},
    {"text": "Break-Even: חודש 8-10", "size": 16, "color": GREEN},
]
add_multiline_box(slide, Inches(0.9), Inches(4.2), Inches(5.2), Inches(2.5), ask_details,
                  alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: Roadmap
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "תוכנית עבודה – 12 חודשים", font_size=36, color=NAVY, bold=True)

quarters = [
    ("Q1: חודשים 1-3", "בנייה", [
        "פיתוח MVP – Scoring Engine",
        "חתימת LOI עם 2 חברות",
        "גיוס 10 מועמדים ראשונים",
        "הסכם עם 3 בתי ספר",
    ], TEAL),
    ("Q2: חודשים 4-6", "פיילוט", [
        "השמת 15 מועמדים בבי\"ס",
        "הפעלת Dual Scoring",
        "איסוף משוב ראשוני",
        "שיפור אלגוריתם",
    ], GOLD),
    ("Q3: חודשים 7-9", "צמיחה", [
        "הרחבה ל-30 משתתפים",
        "Break-Even",
        "חתימת 5 חברות נוספות",
        "פיתוח Feedback Loop",
    ], GREEN),
    ("Q4: חודשים 10-12", "סקייל", [
        "50+ משתתפים פעילים",
        "סיבוב גיוס Seed",
        "הרחבה ארצית",
        "שיתוף פעולה עם צה\"ל",
    ], CORAL),
]

q_w = Inches(2.8)
q_h = Inches(4.5)
q_gap = Inches(0.35)

for i, (title, phase, items, accent) in enumerate(quarters):
    qx = Emu(Inches(0.65).emu + i * (q_w.emu + q_gap.emu))
    qy = Inches(1.5)

    add_rect(slide, qx, qy, q_w, q_h, WHITE)
    add_rect(slide, qx, qy, q_w, Inches(0.08), accent)

    # Phase badge
    badge = add_rect(slide, Emu(qx.emu + q_w.emu // 2 - Inches(0.6).emu), Emu(qy.emu + Inches(0.3).emu),
                     Inches(1.2), Inches(0.4), accent)
    add_text_box(slide, Emu(qx.emu + q_w.emu // 2 - Inches(0.6).emu), Emu(qy.emu + Inches(0.3).emu),
                 Inches(1.2), Inches(0.4),
                 phase, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Emu(qx.emu + Inches(0.2).emu), Emu(qy.emu + Inches(0.9).emu),
                 Emu(q_w.emu - Inches(0.4).emu), Inches(0.5),
                 title, font_size=14, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text_box(slide, Emu(qx.emu + Inches(0.25).emu),
                     Emu(qy.emu + Inches(1.6).emu + Inches(j * 0.55).emu),
                     Emu(q_w.emu - Inches(0.5).emu), Inches(0.45),
                     "◂ " + item, font_size=12, color=GRAY_TEXT)

# Timeline bar at bottom
add_rect(slide, Inches(0.65), Inches(6.3), Inches(12.0), Inches(0.15), RGBColor(0xE2, 0xE8, 0xF0))
for i in range(4):
    seg_w = Inches(3.0)
    add_rect(slide, Emu(Inches(0.65).emu + i * seg_w.emu), Inches(6.3),
             seg_w, Inches(0.15), [TEAL, GOLD, GREEN, CORAL][i])


# ════════════════════════════════════════════════════════════════
# SLIDE 11: Team
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
             "הצוות", font_size=36, color=WHITE, bold=True)

team_members = [
    ("מייסד ומנכ\"ל", "רקע מבצעי-טכנולוגי בצה\"ל\nמנהל רשת + חמ\"ליסט\nחזון ואסטרטגיה", GOLD),
    ("CTO & שותף מייסד", "יזם סדרתי + מתכנת Full-Stack\nניסיון בבניית סטארטאפ\nאחראי על פיתוח ה-MVP", TEAL),
]

member_w = Inches(5.0)
member_h = Inches(4.0)

for i, (role, desc, accent) in enumerate(team_members):
    mx = Emu(Inches(1.5).emu + i * (member_w.emu + Inches(1.0).emu))
    my = Inches(1.8)

    add_rect(slide, mx, my, member_w, member_h, RGBColor(0x1B, 0x25, 0x55))
    add_rect(slide, mx, my, member_w, Inches(0.06), accent)

    # Avatar circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Emu(mx.emu + member_w.emu // 2 - Inches(0.5).emu), Emu(my.emu + Inches(0.4).emu),
        Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()

    icon = "👨‍💼" if i == 0 else "👨‍💻"
    add_text_box(slide, Emu(mx.emu + member_w.emu // 2 - Inches(0.4).emu),
                 Emu(my.emu + Inches(0.55).emu), Inches(0.8), Inches(0.7),
                 icon, font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Emu(mx.emu + Inches(0.3).emu), Emu(my.emu + Inches(1.7).emu),
                 Emu(member_w.emu - Inches(0.6).emu), Inches(0.5),
                 role, font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

    for j, line in enumerate(desc.split("\n")):
        add_text_box(slide, Emu(mx.emu + Inches(0.3).emu),
                     Emu(my.emu + Inches(2.4).emu + Inches(j * 0.4).emu),
                     Emu(member_w.emu - Inches(0.6).emu), Inches(0.4),
                     line, font_size=14, color=RGBColor(0xA0, 0xAE, 0xC0), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: Call to Action
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)
add_rect(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GOLD)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
             "מעוף", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.6),
             "הגשר בין צה\"ל, החינוך, וההייטק", font_size=28, color=TEAL, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), GOLD)

# Key message
add_multiline_box(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(2.0), [
    {"text": "3,500 משרות STEM פנויות  ×  15,000 חיילים עם רקע טכנולוגי", "size": 20, "color": GOLD, "bold": True},
    {"text": "", "size": 10, "color": WHITE},
    {"text": "AI שמחבר בין הצרכים, מדרג, משבץ ומשפר – אוטומטית", "size": 18, "color": WHITE},
], alignment=PP_ALIGN.CENTER)

# CTA buttons
cta_items = [
    ("noam@maof.tech", TEAL),
    ("200K ₪ מענק תנופה", GOLD),
    ("פיילוט 30 חיילים", GREEN),
]

for i, (text, clr) in enumerate(cta_items):
    bx = Emu(Inches(2.0).emu + i * Inches(3.3).emu)
    by = Inches(5.5)

    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        bx, by, Inches(2.8), Inches(0.7))
    btn.fill.solid()
    btn.fill.fore_color.rgb = clr
    btn.line.fill.background()

    add_text_box(slide, bx, Emu(by.emu + Inches(0.12).emu), Inches(2.8), Inches(0.5),
                 text, font_size=16, color=WHITE if clr != GOLD else DARK_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "Confidential | מעוף – Tech-Lead Israel | 2026",
             font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ═══ Save ═══
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "maof-pitch-deck.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
